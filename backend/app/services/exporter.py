"""
uriv-syncboard / backend / app / services / exporter.py
─────────────────────────────────────────────────────────
Converts a persisted DB session into any supported export format.
Returns raw bytes (suitable for FastAPI StreamingResponse).
"""

from __future__ import annotations

import io
import json
import logging
from datetime import datetime
from pathlib import Path
from typing import Literal

import cv2
import numpy as np
from PIL import Image

from app.db.models import Session as DBSession, Note

log = logging.getLogger(__name__)

ExportFmt = Literal["pdf", "docx", "pptx", "markdown", "txt", "json"]


class ExportService:

    def __init__(self, session: DBSession):
        self.session = session
        self.pages   = sorted(session.notes, key=lambda n: n.sequence_order)

    def build(self, fmt: ExportFmt) -> bytes:
        fn = {
            "pdf":      self._pdf,
            "docx":     self._docx,
            "pptx":     self._pptx,
            "markdown": self._markdown,
            "txt":      self._txt,
            "json":     self._json,
        }[fmt]
        buf = io.BytesIO()
        fn(buf)
        return buf.getvalue()

    # ── Helpers ───────────────────────────────────────────────────────────────

    def _decode_image(self, note: Note) -> "np.ndarray | None":
        if not note.image_data:
            return None
        arr = np.frombuffer(note.image_data, dtype=np.uint8)
        return cv2.imdecode(arr, cv2.IMREAD_COLOR)

    def _pil_image(self, note: Note) -> "Image.Image | None":
        img = self._decode_image(note)
        if img is None:
            return None
        return Image.fromarray(cv2.cvtColor(img, cv2.COLOR_BGR2RGB))

    def _header(self) -> str:
        s = self.session
        return (
            f"Session : {s.name}\n"
            f"Date    : {s.started_at.strftime('%Y-%m-%d  %H:%M')}\n"
            f"Pages   : {len(self.pages)}\n"
        )

    # ── PDF ───────────────────────────────────────────────────────────────────

    def _pdf(self, buf: io.BytesIO):
        from reportlab.lib.pagesizes import A4
        from reportlab.platypus import (
            SimpleDocTemplate, Paragraph, Spacer,
            Image as RLImage, PageBreak, HRFlowable,
        )
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import inch
        from reportlab.lib import colors as C

        doc    = SimpleDocTemplate(buf, pagesize=A4,
                                   leftMargin=0.75*inch, rightMargin=0.75*inch,
                                   topMargin=0.75*inch,  bottomMargin=0.75*inch)
        styles = getSampleStyleSheet()
        NAVY   = C.HexColor("#1a1a2e")
        GREY   = C.HexColor("#666666")
        RED    = C.HexColor("#e94560")

        T = ParagraphStyle("T", parent=styles["Title"],   fontSize=24, textColor=NAVY)
        M = ParagraphStyle("M", parent=styles["Normal"],  fontSize=9,  textColor=GREY)
        H = ParagraphStyle("H", parent=styles["Heading2"],fontSize=13, textColor=NAVY)
        B = ParagraphStyle("B", parent=styles["Normal"],  fontSize=10.5, leading=15)

        story = [
            Paragraph(self.session.name, T),
            HRFlowable(width="100%", thickness=2, color=RED, spaceAfter=8),
        ]
        for line in self._header().splitlines():
            story.append(Paragraph(line, M))
        story.append(PageBreak())

        for note in self.pages:
            story.append(Paragraph(f"Page {note.sequence_order}", H))
            story.append(Paragraph(
                f"{note.created_at.strftime('%H:%M:%S')}  ·  Confidence {note.confidence:.1f}%", M
            ))
            pil = self._pil_image(note)
            if pil:
                tmp = io.BytesIO()
                pil.save(tmp, "PNG")
                tmp.seek(0)
                story.append(RLImage(tmp, width=6*inch, height=4*inch, kind="proportional"))
            story.append(Spacer(1, 0.1*inch))
            story.append(Paragraph("Extracted Text", H))
            for line in (note.ocr_text or "").strip().splitlines():
                if line.strip():
                    story.append(Paragraph(line.strip(), B))
            if not (note.ocr_text or "").strip():
                story.append(Paragraph("(No text detected)", M))
            story.append(PageBreak())

        doc.build(story)

    # ── DOCX ──────────────────────────────────────────────────────────────────

    def _docx(self, buf: io.BytesIO):
        from docx import Document
        from docx.shared import Inches, Pt, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH

        doc = Document()
        doc.core_properties.title   = self.session.name
        doc.core_properties.subject = "SyncBoard Export"

        h = doc.add_heading(self.session.name, level=0)
        h.alignment = WD_ALIGN_PARAGRAPH.CENTER

        meta = doc.add_paragraph()
        for line in self._header().splitlines():
            run = meta.add_run(line + "\n")
            run.font.size      = Pt(9)
            run.font.color.rgb = RGBColor(0x88, 0x88, 0x88)
        doc.add_page_break()

        for note in self.pages:
            doc.add_heading(f"Page {note.sequence_order}", level=1)
            sub = doc.add_paragraph()
            r   = sub.add_run(
                f"{note.created_at.strftime('%H:%M:%S')}  ·  Confidence {note.confidence:.1f}%"
            )
            r.font.size      = Pt(9)
            r.font.color.rgb = RGBColor(0x99, 0x99, 0x99)

            pil = self._pil_image(note)
            if pil:
                tmp = io.BytesIO()
                pil.save(tmp, "PNG")
                tmp.seek(0)
                doc.add_picture(tmp, width=Inches(5.5))
                doc.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER

            doc.add_heading("Extracted Text", level=2)
            text = (note.ocr_text or "").strip()
            for line in text.splitlines():
                if line.strip():
                    doc.add_paragraph(line.strip())
            if not text:
                p = doc.add_paragraph("(No text detected)")
                p.runs[0].font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
            doc.add_page_break()

        doc.save(buf)

    # ── PPTX ──────────────────────────────────────────────────────────────────

    def _pptx(self, buf: io.BytesIO):
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor as RGB

        prs = Presentation()
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)

        # Title slide
        sl = prs.slides.add_slide(prs.slide_layouts[0])
        sl.shapes.title.text    = self.session.name
        sl.placeholders[1].text = (
            f"Date: {self.session.started_at.strftime('%B %d, %Y')}\n"
            f"Pages: {len(self.pages)}"
        )

        for note in self.pages:
            sl = prs.slides.add_slide(prs.slide_layouts[6])  # blank

            hdr = sl.shapes.add_textbox(Inches(0.3), Inches(0.1), Inches(12), Inches(0.5))
            tf  = hdr.text_frame
            tf.text = (
                f"Page {note.sequence_order}   ·   "
                f"{note.created_at.strftime('%H:%M:%S')}   ·   "
                f"Confidence {note.confidence:.1f}%"
            )
            run = tf.paragraphs[0].runs[0]
            run.font.size  = Pt(13)
            run.font.bold  = True
            run.font.color.rgb = RGB(0x1a, 0x1a, 0x2e)

            pil = self._pil_image(note)
            if pil:
                tmp = io.BytesIO()
                pil.save(tmp, "PNG")
                tmp.seek(0)
                sl.shapes.add_picture(tmp, Inches(0.3), Inches(0.75),
                                      width=Inches(7.5), height=Inches(6.1))

            tb  = sl.shapes.add_textbox(Inches(8.2), Inches(0.75), Inches(4.8), Inches(6.3))
            tf2 = tb.text_frame
            tf2.word_wrap = True
            tf2.text = (note.ocr_text or "").strip()[:600] or "(No text detected)"
            for para in tf2.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(11)

        prs.save(buf)

    # ── Markdown ──────────────────────────────────────────────────────────────

    def _markdown(self, buf: io.BytesIO):
        lines = [
            f"# {self.session.name}\n",
            f"> **Date:** {self.session.started_at.strftime('%Y-%m-%d %H:%M')}  \n",
            f"> **Pages:** {len(self.pages)}\n\n---\n\n",
        ]
        for note in self.pages:
            lines.append(f"## Page {note.sequence_order}\n\n")
            lines.append(
                f"*{note.created_at.strftime('%H:%M:%S')} — "
                f"confidence {note.confidence:.1f}%*\n\n"
            )
            if note.ocr_text and note.ocr_text.strip():
                lines.append("### Extracted Text\n\n")
                lines.append(note.ocr_text.strip())
                lines.append("\n\n")
            else:
                lines.append("*No text detected.*\n\n")
            lines.append("---\n\n")

        buf.write("".join(lines).encode("utf-8"))

    # ── TXT ──────────────────────────────────────────────────────────────────

    def _txt(self, buf: io.BytesIO):
        lines = ["=" * 60 + "\n", self._header(), "=" * 60 + "\n\n"]
        for note in self.pages:
            lines.append(
                f"── Page {note.sequence_order}  "
                f"[{note.created_at.strftime('%H:%M:%S')}]  "
                f"conf:{note.confidence:.1f}% ──\n\n"
            )
            lines.append((note.ocr_text or "").strip() or "(No text detected)")
            lines.append("\n\n")
        buf.write("".join(lines).encode("utf-8"))

    # ── JSON ─────────────────────────────────────────────────────────────────

    def _json(self, buf: io.BytesIO):
        data = {
            "session_id": str(self.session.id),
            "name":       self.session.name,
            "started_at": self.session.started_at.isoformat(),
            "page_count": len(self.pages),
            "pages": [
                {
                    "sequence_order": n.sequence_order,
                    "created_at":     n.created_at.isoformat(),
                    "ocr_text":       n.ocr_text,
                    "confidence":     round(n.confidence, 2),
                }
                for n in self.pages
            ],
        }
        buf.write(json.dumps(data, indent=2, ensure_ascii=False).encode("utf-8"))
